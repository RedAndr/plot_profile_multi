#!/usr/bin/env python3
# -*- coding: utf-8 -*-
#
# Plot time series for mercury air concentrations
# Andrei Ryjkov, ARQI, ECCC
#

import sys, os, subprocess, argparse, calendar, timeit, shutil
import ray
from ray.util.multiprocessing import Pool
import numpy as np
from datetime import datetime, date, timedelta
from pylab import figure, NaN, clf
import matplotlib as mpl
from matplotlib.dates import YearLocator, MonthLocator, DayLocator, DateFormatter, date2num, num2date
from matplotlib.backends.backend_pdf import PdfPages
import matplotlib.pyplot as plt
from scipy.interpolate import LSQUnivariateSpline       
from scipy.stats import linregress, norm
from PIL import Image
from pdf2image import convert_from_path
from TaylorDiagram import ModTaylorDiagram
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pdfrw import PdfReader, PdfWriter
import rich.traceback
from SolarDiagram import SolarDiagram
from BiasCorrDiagram import bias_correlation_diagram


rich.traceback.install()

def year_range_check(arg):
    """Custom validation function to check if year is within the range (1990-2025)."""
    MIN_YEAR, MAX_YEAR = 1990, 2025
    try:
        f = int(arg)
    except ValueError:    
        raise argparse.ArgumentTypeError("Must be an integer number.")
    if f < MIN_YEAR or f > MAX_YEAR:
        raise argparse.ArgumentTypeError("Argument must be < " + str(MAX_YEAR) + " and > " + str(MIN_YEAR) + ' .')
    return f
  
def month_range_check(arg):
    """Custom validation function to check if month is within the range (1-12)."""
    MIN_MON, MAX_MON = 1, 12
    try:
        f = int(arg)
    except ValueError:    
        raise argparse.ArgumentTypeError("Must be an integer number.")
    if f < MIN_MON or f > MAX_MON:
        raise argparse.ArgumentTypeError("Argument must be < " + str(MAX_MON) + " and > " + str(MIN_MON) + ' .')
    return f

# Script version
ver = '2.71'

# Command line arguments
parser = argparse.ArgumentParser(description="Plot mercury time series")

# Options groups
general_group = parser.add_argument_group("General options")
general_group.add_argument('--version',                        action='version', version=f'%(prog)s {ver}')
general_group.add_argument("-n", "--notrans", dest="trans",    action='store_false',            help="Disable transparency")
general_group.add_argument("-s", "--serial" , dest="parallel", action='store_false',            help="Process in serial, not in parallel")
general_group.add_argument("-d", "--dry-run", dest="dryrun",   action="store_true",             help="Perform a dry run, simulating program execution without any actual changes")
time_group = parser.add_argument_group("Time range options")
time_group.add_argument("year",    type=year_range_check,                                       help="Target year")
time_group.add_argument("--start", type=month_range_check, default=1 ,                          help="Start month (default:  1)")
time_group.add_argument("--end"  , type=month_range_check, default=12,                          help="End month   (default: 12)")
data_group = parser.add_argument_group("Data processing options")
data_group.add_argument("-debias", dest="debias", action='store_true',                          help="Remove bias from the model data")
plot_group = parser.add_argument_group("Plotting options")
plot_group.add_argument("-nodaily",   dest="daily",    action='store_false',                    help="Do not plot daily average")
plot_group.add_argument("-weekly",    dest="weekly",   action='store_true',                     help="Plot weekly average")
plot_group.add_argument("-nomonthly", dest="monthly",  action='store_false',                    help="Do not plot monthly average")
plot_group.add_argument("-nostats",   dest="stats",    action='store_false',                    help="Do not write statistics on the plots")
plot_group.add_argument("-statplot",  dest="statplot", action='store_true',                     help="Plot statistics plots")
plot_group.add_argument("-taylor",    dest="taylor",   action='store_true',                     help="Plot Taylor diagram")
plot_group.add_argument("-solar",     dest="solar",    action='store_true',                     help="Plot Solar diagram")
plot_group.add_argument("-biascor",   dest="biascor",  action='store_true',                     help="Plot Bias-Correlatin diagram")
plot_group.add_argument("-diurnal",   dest="diurplot", type=int, choices=[1,2,3,4], nargs='*',  help="Plot diurnal profiles, can accept monthly and hourly averages")
plot_group.add_argument("-coords",    dest="lox",      action='store_true',                     help="Add sites locations to the legend")
axis_group = parser.add_argument_group("Axis limits options")
axis_group.add_argument("-ymin",  type=float,                                                   help="Y-axis min for time series")
axis_group.add_argument("-ymax",  type=float,                                                   help="Y-axis max for time series")
axis_group.add_argument("-ydmin", type=float,                                                   help="Y-axis min for diurnal plot")
axis_group.add_argument("-ydmax", type=float,                                                   help="Y-axis max for diurnal plot")
output_group = parser.add_argument_group("Output options")
output_group.add_argument("-tz",     type=float, default=0,                                     help="Local time zone (default: 0)")
output_group.add_argument("-hist",   type=int,   default=0,  dest="nbins", nargs='?',           help="Plot histogram with 'NBINS' bins, if 'NBINS' is not specified default is 16")
output_group.add_argument("-pngdpi", type=int,   default=0,  nargs='?',                         help="Create png files with 'PNGDPI' resolution and collect them into 'pptx' file, default: 100")
output_group.add_argument("-pngclr", type=int,   default=64, nargs='?',                         help="Number of colors for png, default: 64")
output_group.add_argument("-pngopt", action='store_true',                                       help="Optimize png images for ppt files")
args = parser.parse_args()

if args.nbins==None:
    hist = 16                                                                            # plot histogram with default number of bins
    histmsg = "histogram with 16 bins"
elif args.nbins==0:
    hist = 0
    histmsg = "no histogram"
else:
    hist = args.nbins                                                                    # plot histogram with specified number of bins
    histmsg = "histogram with %u bins" % hist

if args.pngdpi == None:
    dpi = 100
else:
    dpi = args.pngdpi

if args.pngclr == None:
    pngclr = 64
else:
    pngclr = args.pngclr

if args.diurplot == None:               # to plot diurnal profiles for each month
    diurplot = False
    dim, dih = 0, 0
    diurmsg = f"no diurnal"
else:
    diurplot = True
    if args.diurplot == []:             # default settings for the diurnal plot
        # Diurnal averages
        #dim, dih = 1, 1                # monthly and hourly
        dim, dih = 2, 2                 # bi-monthly and bi-hourly
        #dim, dih = 3, 3                # tri-monthly and tri-hourly
        #dim, dih = 1, 2                # monthly and bi-hourly
    else:
        dim, dih = args.diurplot[0], args.diurplot[1]
    diurmsg = f"diurnal[{dim}:{dih}]"

if args.start > args.end:
    print(f"Error in months range: {args.start} > {args.end}")
    exit(-1)

print("Plotting ", args.year, " year from "+calendar.month_name[args.start]+" to "+calendar.month_name[args.end]+" with "+\
    { True:'daily'  , False:'no daily'  }[args.daily   ]+" and "+\
    { True:'weekly' , False:'no weekly' }[args.weekly  ]+" and "+\
    { True:'monthly', False:'no monthly'}[args.monthly ]+" time series and "+\
    { True:'stats'  , False:'no stats'  }[args.stats   ]+" and "+\
    { True:'stats'  , False:'no stats'  }[args.statplot]+" plots and "+diurmsg+" plots and "+\
    { True:'Taylor' , False:'no Taylor' }[args.taylor  ]+" diagram and "+\
    { True:'BiasCor', False:'no BiasCor'}[args.biascor ]+" diagram and "+\
    { True:'Solar'  , False:'no Solar'  }[args.solar   ]+" diagram and "+ histmsg,
    end=' ')

if args.year == 2013:
  obsrun    =  'BASE'
  sites     = ['LittleFoxLake', 'Kejimkujik', 'Egbert']
  runs      = ['NAVEB93', 'F5B' ]
  runsnames = ['NAVEB93', 'F5B' ]
  colors    = ['Green'  , 'Blue']
elif args.year == 2015:
  obsrun    =  'BASE'
  sites     = ['LittleFoxLake', 'Kejimkujik', 'Egbert']
  runs      = ['NAVEB93', 'F5B' ]
  runsnames = ['NAVEB93', 'F5B' ]
  colors    = ['Green'  , 'Blue']
else:
  print ("Unknown year: ", args.year)
  sys.exit(-1)

# Time zone correction for some observations
#TZ = { 'Alert':-5, 'Nord':0, 'Zeppelin':1, 'Andoya':1, 'Denali':-9, 'LittleFoxLake':-8,
       #'Lahemaa':2, 'Birkenes':1, 'FortMcKay':-7, 'FortMcMurray':-7,
       #'Pensacola':-6, 'Beltsville':-5, 'Underhill':-5, 'Stilwell':-6, 'Grand_Bay':-6,
       #'Horicon_Marsh':-6, 'Yorkville':-6, 'Rochester':-5, 'Bronx':-5,'Salt_Lake_City':-7,
       #'Huntington':-5, 'Athens':-5, 'Piney_Reservoir':-5, 'South_Bass_Island':-5, 'Birmingham':-6,
       #'MaunaLoa':-10}
TZ = { 'Denali':-9, 'Pensacola':-6, 'Beltsville':-5, 'Underhill':-5, 'Stilwell':-6, 'Grand_Bay':-6,
       'Horicon_Marsh':-6, 'Yorkville':-6, 'Rochester':-5, 'Bronx':-5,'Salt_Lake_City':-7,
       'Huntington':-5, 'Athens':-5, 'Piney_Reservoir':-5, 'South_Bass_Island':-5, 'Birmingham':-6,
       'MaunaLoa':-10}

# Sites locations
Lox = {'Alert':(82.50,-62.33), 'Nord':(81.60,-16.67), 'Zeppelin':(78.90,11.88), 'Andoya':(69.28,16.01), 'Pallas':(68.00,24.24), 'Bredkalen':(63.85,15.33), 'LittleFoxLake':(61.35,-135.63), 'Lahemaa':(59.50,25.90), 'BirkenesII':(58.38,8.25), 'Birkenes':(58.38,8.25), 'Rao':(57.39,11.91), 
       'LowerCamp':(57.03,-111.50), 'FortMcMurray':(56.75,-111.48), 'Vavihill':(56.02,13.15), 'Auchencorth_Moss':(55.79,-3.24), 'FlinFlon':(54.77,-101.88), 'Zingst':(54.43,12.73), 'DiablaGora':(54.15,22.07), 'MaceHead':(53.33,-9.90), 'Waldhof':(52.80,10.75), 'Listvyanka':(51.85,104.89), 
       'Schmucke':(50.65,10.78), 'Mingan':(50.27,-64.23), 'Whistler':(50.07,-122.93), 'Kosetice':(49.57,15.08), 'Ucluelet':(48.92,-125.54), 'Saturna':(48.78,-123.13), 'Schauinsland':(47.91,7.81), 'PresqueIsle':(46.70,-68.03), 'Iskrba':(45.56,14.86), 'St.Anicet':(45.12,-74.28), 
       'Halifax':(44.67,-63.57), 'Underhill':(44.53,-72.87), 'Kejimkujik':(44.43,-65.21), 'Egbert':(44.23,-79.78), 'Huntington':(43.97,-74.22), 'HoriconMarsh':(43.47,-88.62), 'Rochester':(43.15,-77.55), 'MtChangbai':(42.40,128.11), 'SouthBassIsland':(41.66,-82.83), 'Bronx':(40.87,-73.88), 
       'SaltLakeCity':(40.71,-111.96), 'Piney_Reservoir':(39.71,-79.01), 'Longobucco':(39.39,16.61), 'Athens':(39.31,-82.12), 'Beltsville':(39.03,-76.82), 'MtWalinguan':(36.29,100.90), 'Stilwell':(35.75,-94.67), 'Yorkville':(33.93,-85.05), 'Birmingham':(33.52,-86.81), 
       'Pensacola':(30.43,-87.20), 'Grand_Bay':(30.43,-88.43), 'MaunaLoa':(19.54,-155.58), 'Manaus':(-2.89,-59.97), 'CapePoint':(-34.35,18.49), 'Bariloche':(-41.10,-71.40), 'AmsterdamIsland':(-37.80,77.55), 'Dumont_d\'Urville':(-66.66,140.01), 'Troll':(-72.02,2.53), 'Concordia':(-75.10,123.35),
       'Denali':(63.72,-148.97), 'LittleFoxLake':(61.35,-135.63), 'FortMcKay':(57.15,-111.64), 'LowerCamp':(57.03,-111.50), 'FortMcMurray':(56.75,-111.48), 'Mingan':(50.27,-64.23), 'Whistler':(50.05,-122.95), 'Ucluelet':(48.92,-125.54), 'Saturna':(48.78,-123.13), 'PresqueIsle':(46.70,-68.03), 
       'St-Anicet':(45.12,-74.28), 'Underhill':(44.53,-72.87), 'Kejimkujik':(44.43,-65.21), 'Egbert':(44.23,-79.78), 
       'Huntington':(43.97,-74.22), 'Horicon_Marsh':(43.47,-88.62), 'RochesterB':(43.15,-77.55), 'South_Bass_Island':(41.66,-82.83), 'Bronx':(40.87,-73.88), 'Salt_Lake_City':(40.71,-111.96), 'Piney_Reservoir':(39.70,-79.01), 
       'Athens':(39.31,-82.12), 'Beltsville':(39.03,-76.82), 'Stilwell':(35.75,-94.67), 'Yorkville':(33.93,-85.05), 'Birmingham':(33.52,-86.81), 'Pensacola':(30.43,-87.20), 'Grand_Bay_NERR':(30.43,-88.43), 
       'Pallas':(68.00,24.24), 'Bredkalen':(63.85,15.33), 'Lahemaa':(59.50,25.90), 'Birkenes':(58.38,8.25), 'Rao':(57.39,11.91), 'Vavihill':(56.02,13.15), 'Auchencorth_Moss':(55.79,-3.24), 'Zingst':(54.43,12.73), 'DiablaGora':(54.15,22.07), 'MaceHead':(53.33,-9.90), 'Waldhof':(52.80,10.75), 
       'Schmucke':(50.65,10.78), 'Kosetice':(49.57,15.08), 'Schauinsland':(47.91,7.81), 'Iskrba':(45.56,14.86), 
       'MaunaLoa':(19.54,-155.58), 'CapePoint':(-34.35,18.49), 'AmsterdamIsland':(-37.80,77.55), '':(0.00,0.00), 'Alert':(82.50,-62.33), 'Nord':(81.60,-16.67), 'Zeppelin':(78.90,11.88), 'Andoya':(69.28,16.01), 'Dumont_d\'Urville':(-66.66,140.01), 'Troll':(-72.02,2.53),
       'ATARS':(-12.249,131.045), 'Kodaikanal':(10.2317,77.46524), 'CapeHedo':(26.8643,128.25141), 'Sisal':(21.16356,-90.04679), 'Toolik':(68.63,-149.6), 'Harwell':(51.573056,-1.316667)}

print ("for the following runs: ", runs, " and the following sites: ", sites)

if dpi>0:
    print(f"Create PPTX file with PNG images [dpi:{dpi}, colors:{pngclr}, pngopt:{args.pngopt}].")
    print(f"PngOpt: {shutil.which('optipng')} {shutil.which('pngout')} {shutil.which('advpng')}")

if args.dryrun:
    print('Dry run finished.')
    exit()

plt.ioff()
mpl.use('Agg')
mpl.pyplot.style.use('ggplot')                          # Plotting style for matplotlib
np.seterr(divide='ignore')

metadata = {'Title':'Mercury model validation',         # metadata for pdf
            'Author':'Andrei Ryjkov ARQI ECCC',
            'Creator':f'plot_profile_multi {ver}'}
mpl.rcParams['pdf.compression'] = 1                     # Compression level for pdf in matplotlib, isn't necessary if pdfrw_compress=True or cpdf=True
cpdf = shutil.which('cpdf')                             # Squeeze pdf using cpdf tool if found
if cpdf:
    pdfrw_compress = False
else:
    pdfrw_compress = True

is_leap_year = args.year%4 == 0 and (args.year%100 != 0 or args.year%400 == 0)
days_in_year = 366 if is_leap_year else 365

date0   = int(date2num(datetime.strptime('01/01/'+"%4u"%args.year+' 00:00', '%d/%m/%Y %H:%M'))) - 1     # First day of the year minus one day
datemin = date(args.year, args.start, 1                                         )                       # start time - first days of beggining month
datemax = date(args.year, args.end  , calendar.monthrange(args.year,args.end)[1])                       # end time - last day of end month
datemsk = np.array(np.arange(datemin,datemax),bool)
datemsk = np.append(datemsk, np.array(np.zeros((days_in_year+1)-datemsk.shape[0]),bool))                # Date mask for daily data, works only if start month = 1

date1 = date(args.year-1, 12, 31)                                                                       # 1st  day of the year - Jan  1 - 1 day
date2 = date(args.year+1,  1,  1)                                                                       # Last day of the year - Dec 31 + 1 day
datemskhr1 = np.array(np.arange(date1  , datemin, np.timedelta64(1,'h')), bool)
datemskhr2 = np.array(np.arange(datemin, datemax, np.timedelta64(1,'h')), bool)
datemskhr3 = np.array(np.arange(datemax, date2  , np.timedelta64(1,'h')), bool)
datemskhr1[:] = False
datemskhr2[:] = True
datemskhr3[:] = False
datemskhr = np.array(np.append(np.append(datemskhr1, datemskhr2), datemskhr3), bool)


def smooth(x, signal, window_size):
    """
    Smooths a signal using either spline interpolation or convolution with a Hamming window.

    Args:
        x (numpy.ndarray): Array of independent variable values (e.g., time points).
        signal (numpy.ndarray): Array of signal values to be smoothed.
        window_size (int): Size of the window for convolution (used if positive).

    Returns:
        numpy.ndarray: Smoothed signal array.
    """

    if window_size < 0:                                                                                   # Use spline interpolation for negative window_size
        # Removed UnivariateSpline as LSQUnivariateSpline is generally preferred
        t = np.linspace(x[1], x[-2], num=24)                                                                # Create new evaluation points for spline
        spl = LSQUnivariateSpline(x, signal, t)                                                             # Fit a spline to the signal
        y = spl(x)                                                                                          # Evaluate the spline at original x values
        return y
    else:                                                                                                 # Use convolution with Hamming window for positive window_size
        # Pad the signal with mirrored edges for boundary handling
        s = np.concatenate((signal[window_size:0:-1], signal, signal[-1:-window_size:-1]))

        # Use Hamming window for smoother filtering
        w = np.hamming(window_size)

        # Handle missing values (NaNs) before convolution
        msk = np.isnan(s)                                                                                   # Create mask for missing values
        s[msk] = np.interp(np.flatnonzero(msk), np.flatnonzero(~msk), s[~msk])                              # Linear interpolation

        # Perform convolution for smoothing
        y = np.convolve(w / w.sum(), s, mode="same")                                                        # Normalize window weights

        # Restore NaNs after convolution
        y[msk] = np.nan

        # Return the central part of the smoothed signal (excluding padded edges)
        return y[window_size:-window_size + 1]


def average_daily(obs_dat, obs_val):
    """
    Calculates the daily average of observation values, filling missing days with NaNs.

    Args:
        obs_dat (numpy.ndarray): Array of observation dates (e.g., timestamps).
        obs_val (numpy.ndarray): Array of corresponding observation values.

    Returns:
        tuple: Tuple containing two NumPy arrays:
            - obs_dat_day (numpy.ndarray): Array of daily dates.
            - obs_val_day (numpy.ndarray): Array of daily average values (or NaNs for missing days).
    """

    # Create arrays to store daily data (including one extra day for edge cases)
    obs_dat_day = np.arange(date0, date0 + days_in_year + 1, dtype=np.int64)
    obs_val_day = np.zeros(days_in_year + 1, dtype=np.float64)
    obs_cnt_day = np.zeros(days_in_year + 1, dtype=np.int64)

    # Iterate through observation data and accumulate values/counts for each day
    for i, day in enumerate(np.trunc(obs_dat - date0).astype(np.int64)):
        try:
            # Handle potential index errors by checking if day is within bounds
            if 0 <= day < len(obs_val_day):  
                obs_val_day[day] += obs_val[i]
                obs_cnt_day[day] += 1
        except IndexError:
            pass  # Silently ignore index errors (out-of-bounds)

    # Handle days with no observations (count as -1 for division later)
    obs_cnt_day[obs_cnt_day == 0] = -1

    # Calculate daily averages (replacing -1 counts with NaN for missing days)
    obs_val_day = np.divide(obs_val_day, obs_cnt_day, where=obs_cnt_day != -1)
    obs_val_day[obs_cnt_day == -1] = np.nan

    return obs_dat_day, obs_val_day


def average_hourly(observation_dates, observation_values):
    """
    Calculates the hourly average of observation values, filling missing hours with NaNs.

    Args:
        observation_dates (numpy.ndarray): Array of observation dates (e.g., timestamps).
        observation_values (numpy.ndarray): Array of corresponding observation values.

    Returns:
        tuple: Tuple containing two NumPy arrays:
            - hourly_observation_dates (numpy.ndarray): Array of hourly observation timestamps.
            - hourly_average_values (numpy.ndarray): Array of hourly average values (or NaNs for missing hours).
    """

    # Create arrays to store hourly data (including one extra hour for edge cases)
    num_hours = (days_in_year + 1) * 24
    hourly_observation_dates = np.array([date0 + hour / 24. for hour in range(num_hours)], dtype=np.float64)
    hourly_average_values = np.zeros(hourly_observation_dates.shape[0], dtype=np.float64)
    hourly_observation_counts = np.zeros(hourly_observation_dates.shape[0], dtype=np.int64)

    # Iterate through observation data and accumulate values/counts for each hour
    for i, observation_hour in enumerate(np.round((observation_dates - date0) * 24).astype(np.int64)):
        try:
            hourly_average_values[observation_hour] += observation_values[i]
            hourly_observation_counts[observation_hour] += 1
        except:
            pass  # Handle potential errors (e.g., out-of-bounds index)

    # Handle hours with no observations (count as -1 for division later)
    hourly_observation_counts[hourly_observation_counts == 0] = -1

    # Calculate hourly averages (replacing -1 counts with NaN for missing hours)
    hourly_average_values = np.divide(hourly_average_values, hourly_observation_counts, where=hourly_observation_counts != -1)
    hourly_average_values[hourly_observation_counts == -1] = np.nan

    return hourly_observation_dates, hourly_average_values


def read_data(year, obsrun, site, suff):
    """
    Reads data from an obs or mod file, performs daily averaging, missing value handling,
    and smoothing, returning various data arrays.

    Args:
        year (int): Year of the data.
        obsrun (str): Observation run identifier.
        site (str): Site name.
        suff (str): File suffix (e.g., 'M' or 'D').

    Returns:
        list: List containing various data arrays:
            - obs_dat_day (numpy.ndarray): Array of daily observation dates.
            - obs_val_day (numpy.ndarray): Array of daily average values (or NaNs for missing days).
            - smtw_obs_val_day (numpy.ndarray): Array of smoothed daily values (weekly smoothing).
            - smtm_obs_val_day (numpy.ndarray): Array of smoothed daily values (monthly smoothing).
            - diu (numpy.ndarray): Array of diurnal profiles (if `diurplot` is True, otherwise False).
            - obs_dat_hr (numpy.ndarray): Array of hourly observation dates.
            - obs_val_hr (numpy.ndarray): Array of hourly average values (or NaNs for missing hours).
    """

    # Construct filename based on input parameters
    file_name = "%4u.%s/%s%4u_%1s.dat" % (year, obsrun, site, year, suff)

    try:
        # Open the data file (focused try block)
        with open(file_name) as file:
            obs_dat, obs_val = [], []  # Lists to store data points

            # Read lines from the file
            for line in file:
                # Skip empty lines
                if not line.strip():
                    continue

                try:
                    # Split line and convert date/time
                    ls = line.split()
                    newdate = date2num(datetime.strptime(ls[0] + ' ' + ls[1], '%d/%m/%Y %H:%M'))
                except ValueError:
                    print('Error converting [' + ls[0] + ' ' + ls[1] + '] from file: ' + file_name)
                    print('line: [' + line + ']')
                    continue

                try:
                    # Convert value to float
                    val = float(ls[2])
                except ValueError:
                    print('Error converting [' + ls[2] + '] from file: ' + file_name)
                    print('line: [' + line + ']')
                    continue

                # Filter out invalid values (-2.0 to 100.0)
                if val < -2.0 or val > 100.0:
                    continue

                # Append data points (only if new date or empty list)
                if len(obs_dat) == 0 or newdate > obs_dat[-1]:
                    obs_dat.append(newdate)
                    obs_val.append(val)

    except FileNotFoundError:
        print(f"Error: File '{file_name}' not found.")
        sys.exit(-1)

    # Convert lists to NumPy arrays for efficiency
    obs_dat = np.array(obs_dat)
    obs_val = np.array(obs_val)

    # Handle local time to UTC conversion for specific sites (if applicable)
    if suff == 'M' and site in TZ:
        obs_dat -= TZ[site] / 24.

    # Check for empty data
    if len(obs_dat) == 0 or len(obs_val) == 0:
        print("Error: no data in " + file_name)
        sys.exit(-1)

    # Calculate daily averages (using average_daily function)
    obs_dat_day, obs_val_day = average_daily(obs_dat, obs_val)

    # Calculate hourly averages (using average_hourly function)
    obs_dat_hr, obs_val_hr = average_hourly(obs_dat, obs_val)

    # Check for sufficient data for daily averages
    if np.count_nonzero(~np.isnan(obs_val_day)) == 0:
        print(f"Error: not enough data to calculate daily averages in {file_name}")
        sys.exit(-1)

    # Smooth daily data (weekly and monthly) using smooth function
    smtw_obs_val_day = smooth(obs_dat_day, obs_val_day,  7)                             # Weekly smoothing window
    smtm_obs_val_day = smooth(obs_dat_day, obs_val_day, 30)                             # Monthly smoothing window (assuming 30 days)

    # Calculates diurnal profiles for observations if requested
    if diurplot:
        diu = np.zeros([12//dim, 24//dih+1])
        cnt = np.zeros([12//dim, 24//dih+1])
        for i, dt in enumerate(num2date(obs_dat)):
            if dt.year == args.year:
                mn = np.int64((dt.month+(dim-1))/dim-1)
                hr = np.int64(np.trunc(((dt.hour+args.tz)%24)/dih))
                diu[mn, hr] += obs_val[i]
                cnt[mn, hr] += 1
        cnt[cnt==0] = -1
        diu = np.divide(diu, cnt)
        diu[diu==0.0] = NaN                                                             # insert NaNs for gaps between daily profiles
        diu[cnt<5] = NaN                                                                # insert NaNs if not enough points for averaging
    else:
        diu = False

    obs_avg = np.mean(obs_val[np.logical_and(np.logical_and(obs_val>0, obs_val<10), ~np.isnan(obs_val))])
    print ( f"{file_name:50s} {len(obs_dat):6n} points => {len(obs_val_day[~np.isnan(obs_val_day)]):3n} days  Mean: {obs_avg:8.4f}" )

    # (Assuming other calculations using obs_dat_day, obs_val_day, etc.)
    return [obs_dat_day, obs_val_day, smtw_obs_val_day, smtm_obs_val_day, diu, obs_dat_hr, obs_val_hr]


def read_site(site):
    """
    Reads observation and model data for a specific site.

    Args:
        site (str): Name of the site.
        args (argparse.Namespace): Namespace containing command-line arguments
            (including 'year' and 'obsrun').
        runs (list): List of model run identifiers.
        debias (bool, optional): Flag indicating whether to debias model data.
            Defaults to False.

    Returns:
        tuple: Tuple containing various data arrays for the site:
            - site (str): Name of the site.
            - obs_dat (numpy.ndarray): Array of daily observation dates.
            - obs_val (numpy.ndarray): Array of daily observation values (or NaNs for missing days).
            - obs_dat_hr (numpy.ndarray): Array of hourly observation dates.
            - obs_val_hr (numpy.ndarray): Array of hourly observation values (or NaNs for missing hours).
            - smtw_obs_val (numpy.ndarray): Array of smoothed daily values (weekly smoothing).
            - smtm_obs_val (numpy.ndarray): Array of smoothed daily values (monthly smoothing).
            - obs_diu (numpy.ndarray): Array of diurnal profiles (if `diurplot` is True, otherwise False).
            - mods (list): List of model data (each element is a list returned by `read_data`).
    """

    # Read observation data
    obs_dat, obs_val, smtw_obs_val, smtm_obs_val, obs_diu, obs_dat_hr, obs_val_hr = read_data(args.year, obsrun, site, 'M')

    # Handle debiasing for observation data (if applicable)
    if args.debias:
        obs_avg = np.nanmean(obs_val)                   # Calculate mean, ignoring NaNs
        obs_val -= obs_avg                              # Remove bias from observation values

    # Initialize list to store model data
    mods = []

    # Read and debias if requested model data for all runs
    for run in runs:
        mm = read_data(args.year, run, site, 'R')
        if args.debias:
            mod_avg = np.mean(mm[2])                    # Calculate mean of the third element (assuming daily values)
            mm[1] -= mod_avg - obs_avg                  # Debias daily values
            mm[2] -= mod_avg - obs_avg                  # Debias other relevant data (assuming indices)
            mm[3] -= mod_avg - obs_avg                  # Debias additional data (assuming indices)
        mods.append(mm)

    return site, obs_dat, obs_val, obs_dat_hr, obs_val_hr, smtw_obs_val, smtm_obs_val, obs_diu, mods


def RMSE(predictions, targets):
    """
    Calculates the Root Mean Squared Error (RMSE) between predictions and targets.

    Args:
        predictions (numpy.ndarray): Array of predicted values.
        targets (numpy.ndarray): Array of target values.

    Returns:
        float: The root mean squared error.
    """

    # Calculate squared errors
    squared_errors = np.square(predictions - targets)

    # Return RMSE (square root of mean squared error)
    return np.sqrt(squared_errors.mean())


def URMSE(predictions, targets):
    """
    Calculates the unbiased Root Mean Squared Error (URMSE) between predictions and targets.

    Args:
        predictions (numpy.ndarray): Array of predicted values.
        targets (numpy.ndarray): Array of target values.

    Returns:
        float: The unbiased root mean squared error.
    """

    # Calculate bias as the difference between mean predictions and mean targets
    bias = predictions.mean() - targets.mean()

    # Calculate squared errors with bias correction
    squared_errors = np.square((predictions - bias) - targets)

    # Return URMSE (square root of mean squared error)
    return np.sqrt(squared_errors.mean())


def NSTD(predictions, targets):
    """
    Calculates the Normalized Standard Deviation (NSTD) between predictions and targets.

    Args:
        predictions (numpy.ndarray): Array of predicted values.
        targets (numpy.ndarray): Array of target values.

    Returns:
        float: The normalized standard deviation.
    """

    # Calculate standard deviations of predictions and targets
    std_pred = np.std(predictions)
    std_targets = np.std(targets)

    # Return NSTD (ratio of standard deviations)
    return std_pred / std_targets if std_targets != 0 else np.nan  # Handle division by zero


def FAC2(predictions, targets):
    """
    Calculates the Fraction of Agreement within a Factor of 2 (FAC2).

    Args:
        predictions (numpy.ndarray): Array of predicted values.
        targets (numpy.ndarray): Array of target values.

    Returns:
        float: The fraction of agreements within a factor of 2.
    """

    # Calculate ratio of predictions to targets
    predivtar = predictions / targets

    # Count elements within the range of 0.5 to 2.0 (inclusive)
    count_fac2 = np.count_nonzero((0.5 <= predivtar) & (predivtar <= 2.0))

    # Return FAC2 (fraction of elements within the factor of 2 range)
    return count_fac2 / len(targets)


def MFB(predictions, targets):
    """
    Calculates the Mean Fractional Bias (MFB).

    Args:
        predictions (numpy.ndarray): Array of predicted values.
        targets (numpy.ndarray): Array of target values.

    Returns:
        float: The mean fractional bias.
    """

    # Calculate fractional bias for each element
    fractional_bias = (predictions - targets) / (predictions + targets)

    # Return MFB (average of fractional bias)
    return 2.0 * fractional_bias.sum() / len(targets)


class DiurnalPlot:
  """
  A class to represent and plot diurnal data.
  """

  def __init__(self, figsize=(9.5, 5), main_axes_size=None):
    """
    Initializes a DiurnalPlot object.

    Args:
        figsize (tuple, optional): Size of the figure in inches. Defaults to (9.5, 5).
        main_axes_size (tuple, optional): Size of the main axes within the figure. Defaults to None (uses figure size).
    """
    self.fig, self.ax = plt.subplots(figsize=figsize)
    if main_axes_size:
      self.ax.set_aspect('equal')  # Optional: Set equal aspect ratio for axes
      self.ax.patch.set_alpha(0)  # Optional: Set transparent background for axes
      self.ax.set_position(main_axes_size)

  def calculate_x_ticks_and_labels(self, data, data_per_hour, display_dimension):
    """
    Calculates tick positions and labels for the x-axis based on data and display preferences.

    Args:
        data (numpy.ndarray): The data array with hourly values.
        data_per_hour (int): The number of data points per hour in the data.
        display_dimension (int): The desired display dimension (1: monthly, 2: bi-monthly, 3: tri-monthly).

    Returns:
        tuple: A tuple containing lists of x-axis tick positions and labels.
    """
    num_data_points = len(data.flatten())
    ticks = [i * (24 // data_per_hour + 1) for i in range(12//dim+1)]
    if display_dimension == 1:
      labels = [calendar.month_abbr[i  +1] for i in range(12)]
    elif display_dimension == 2:
      labels = [calendar.month_abbr[i*2+1] + '-' + calendar.month_abbr[i * 2 + 2] for i in range(6)]
    elif display_dimension == 3:
      labels = [calendar.month_abbr[i*3+1] + '-' + calendar.month_abbr[i * 3 + 3] for i in range(4)]
    elif display_dimension == 4:
      labels = [calendar.month_abbr[i*4+1] + '-' + calendar.month_abbr[i * 4 + 4] for i in range(3)]
    else:
      raise ValueError("Invalid display dimension. Please choose 1 (monthly), 2 (bi-monthly), 3 (tri-monthly), or 4 (quart-monthly).")
    return ticks, labels + ['']  # Add an empty label at the end

  def plot_data(self, data, data_per_hour, display_dimension, label, color, markersize=6, linewidth=2, alpha=0.85):
    """
    Plots a line representing the provided data on the existing axes.

    Args:
        data (numpy.ndarray): The data array with hourly values.
        data_per_hour (int): The number of data points per hour in the data.
        display_dimension (int): The desired display dimension (1: monthly, 2: bi-monthly, 3: tri-monthly).
        label (str): The label for the data line in the legend.
        color (str): The color of the data line.
        markersize (int, optional): Size of the markers on the line. Defaults to 7.
        linewidth (int, optional): Width of the line. Defaults to 2.
        alpha (float, optional): Transparency of the line (0: fully transparent, 1: fully opaque). Defaults to 0.85.
    """
    ticks, labels = self.calculate_x_ticks_and_labels(data, data_per_hour, display_dimension)
    self.ax.plot(range(len(data.flatten())), data.flatten(), '.-', color=color, antialiased=True, alpha=alpha, markersize=markersize, lw=linewidth, label=label)
    self.ax.set_xticks(ticks)
    self.ax.set_xticklabels(labels, ha='left')

  def finalize_plot(self, y_min=None, y_max=None, title=None, legend_args=None, filename=None):
    """
    Finalizes the plot by setting axes limits, adding title, legend, and saving the figure.

    Args:
        y_min (float, optional): Minimum value for the y-axis. Defaults to None (autoscale).
        y_max (float, optional): Maximum value for the y-axis. Defaults to None (autoscale).
        title (str, optional): Title for the plot. Defaults to None.
        legend_args (dict, optional): A dictionary containing keyword arguments for the legend. Defaults to None.
        filename (str, optional): Filename to save the plot image. Defaults to None (no saving).
    """
    if y_min is None or y_max is None:
      self.ax.autoscale_view()
    else:
      self.ax.set_ylim(y_min, y_max)
    self.ax.grid(True)

    if title:
      self.fig.suptitle(title, fontsize=12)  # Set title using suptitle for better placement

    if legend_args:
      self.ax.legend(**legend_args)  # Unpack legend arguments from the dictionary

    if filename:
      self.fig.savefig(filename)

    return self.fig


class ScatterPlotWithRegression:
  """
  A class to create scatter plots with linear regressions.

  This class allows you to add observed data (obs_val) and corresponding values (days)
  to a plot. It calculates and displays a linear regression line (2-parameter) and 
  a simple linear fit (1-parameter) along with the y=x line.
  """

  def __init__(self, figsize=(9, 8), axes_position=(0.07, 0.07, 0.88, 0.85), xlabel='Obs', ylabel='Mod'):
    """
    Initializes a ScatterPlotWithRegression object.

    Args:
        figsize (tuple, optional): Size of the figure in inches. Defaults to (9, 8).
        axes_position (tuple, optional): Position of the axes within the figure (normalized coordinates). Defaults to (0.07, 0.07, 0.88, 0.85).
        xlabel (str, optional): Label for the x-axis. Defaults to 'Obs'.
        ylabel (str, optional): Label for the y-axis. Defaults to 'Mod'.
    """
    import matplotlib.pyplot as plt
    self.fig, self.ax = plt.subplots(figsize=figsize)
    self.ax.set_xlabel(xlabel)
    self.ax.set_ylabel(ylabel)
    self.ax.set_position(axes_position)
    self.lines = []  # List to store line plot objects
    self.min_data =  0.
    self.max_data = 10.
    self.ax.plot([self.min_data,self.max_data], [self.min_data,self.max_data], '--', color='white', alpha=0.8, lw=1, antialiased=True)      # Plot y=x line (dashed white)

  def add_data(self, observed_values, corresponding_values, color='red', markersize=6, alpha=0.85, antialiased=True):
    """
    Adds a scatter plot with regressions to the existing axes.

    Args:
        observed_values (list): List of observed values.
        corresponding_values (list): List of corresponding values.
        color (str, optional): Color for the plot lines. Defaults to 'red'.
        markersize (int, optional): Size of the markers for the scatter plot. Defaults to 6.
        alpha (float, optional): Transparency of the lines (0: fully transparent, 1: fully opaque). Defaults to 0.85.
        antialiased (bool, optional): Enables antialiasing for smoother lines. Defaults to True.
    """
    # Plot scatter points
    #self.lines.append(
    self.ax.plot(observed_values, corresponding_values, '.', color=color, markersize=markersize, alpha=alpha, antialiased=antialiased)
    #[0])

    # Calculate data extent (min & max values)
    self.min_data = max(self.min_data, min(min(observed_values), min(corresponding_values)))
    self.max_data = min(self.max_data, max(max(observed_values), max(corresponding_values)))

    # Perform linear regression (2-parameter)
    slope, intercept, _, _, _ = linregress(observed_values, corresponding_values)
    regression_line_x = np.linspace(self.min_data, self.max_data, 10)  # Generate x-values for regression line
    regression_line_y = slope * regression_line_x + intercept
    self.lines.append(self.ax.plot(regression_line_x, regression_line_y, '-', color=color, alpha=0.5, lw=3, antialiased=antialiased)[0])

    # Perform simple linear fit (1-parameter)
    fit_coeff = np.linalg.lstsq(observed_values[:, np.newaxis], corresponding_values, rcond=None)[0]
    fit_line_x = np.linspace(self.min_data, self.max_data, 10)
    fit_line_y = fit_coeff * fit_line_x
    #self.lines.append(
    self.ax.plot(fit_line_x, fit_line_y, '-', color=color, alpha=0.2, lw=2, antialiased=antialiased)    #[0])

  def finalize(self, run_names):
    """
    Finalizes the plot by adding a grid, legend, and returning the figure.

    Args:
        run_names (list): List of names for the data sets in the plot.

    Returns:
        matplotlib.figure.Figure: The created plot figure.
    """
    self.min_data = float(int(self.min_data * 10) - 1) / 10
    self.max_data = float(int(self.max_data * 10) + 1) / 10
    # Set axes limits
    self.ax.set_xlim(self.min_data, self.max_data)
    self.ax.set_ylim(self.min_data, self.max_data)
    self.ax.grid(True)
    # Create legend with proper positioning and formatting
    self.ax.legend(self.lines, run_names, bbox_to_anchor=(0, 1.02, 1, 0.1), loc=3, ncol=7, mode="expand", borderaxespad=0)
    return self.fig


def calculate_statistics(observed_hourly_values, hourly_model_values, date_mask_hourly, statistic_dictionary, run_name):
  """
  Calculates statistics for observed and modeled hourly values.

  This function calculates various statistics including slope, intercept, correlation,
  bias, RMSE (Root Mean Squared Error), URMSE (Unbiased RMSE), NRMSE (Normalized RMSE),
  FAC2 (Fraction of explained variance), and MFB (Mean Fractional Bias). It also calculates
  a composite metric (AQPI) combining these statistics. Statistics are stored in the provided
  dictionary with appropriate keys.

  Args:
      observed_hourly_values (np.array): Array of observed hourly values.
      hourly_model_values (np.array): Array of modeled hourly values.
      date_mask_hourly (np.array): Array indicating valid data points (True) and missing values (False).
      statistic_dictionary (dict): Dictionary to store calculated statistics.
      run_name (str): Name of the model run.

  Returns:
      str: Formatted string containing calculated statistics for the current run.
  """

  # Remove missing values (NaN) and points masked by date mask
  valid_data = np.logical_and(np.logical_not(np.isnan(observed_hourly_values)), np.logical_and(np.logical_not(np.isnan(hourly_model_values)), date_mask_hourly))
  filtered_observed_values = observed_hourly_values[valid_data]
  filtered_model_values = hourly_model_values[valid_data]

  # Check if we have any valid data for calculations
  if len(filtered_observed_values) > 0:
    # Perform linear regression
    slope, intercept, correlation, p_value, std_err = linregress(filtered_observed_values, filtered_model_values)

    # Calculate statistics
    model_average = filtered_model_values.mean()
    bias = model_average - filtered_observed_values.mean()
    rmse  = RMSE (filtered_model_values, filtered_observed_values)
    urmse = URMSE(filtered_model_values, filtered_observed_values)
    nrsd  = NSTD (filtered_model_values, filtered_observed_values)
    fac2  = FAC2 (filtered_model_values, filtered_observed_values)
    mfbs  = MFB  (filtered_model_values, filtered_observed_values)

    # Calculate composite metric (AQPI)
    aqpi = (fac2 + correlation + (1.0 - np.abs(0.5 * mfbs))) / 3.0

  else:
    # Set all statistics to zero if no valid data
    slope, intercept, correlation, p_value, standard_error, bias, rmse, urmse, nrsd, fac2, mfbs, aqpi = 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0, 0

  # Format statistics string using f-string
  statistics_string = f"{run_name:<8s}: $\\bar{{x}}_{{mod}}={model_average:.2f}, bias={bias:+6.3f}, R={correlation:.3f}, x_{{mod}}={slope:.2f}*x_{{obs}}{intercept:+5.2f}, \sigma_e={std_err:.3f}, N={len(filtered_observed_values):3g}, RMSE={rmse:.3f}, URMSE={urmse:.3f}$"

  # Update statistic dictionary with keys using f-strings
  statistic_dictionary.update({f"bias {run_name}": abs(bias)})
  statistic_dictionary.update({f"corr {run_name}": correlation})
  statistic_dictionary.update({f"slop {run_name}": slope})
  statistic_dictionary.update({f"intr {run_name}": intercept})
  statistic_dictionary.update({f"stdr {run_name}": std_err})
  statistic_dictionary.update({f"rmse {run_name}": rmse})
  statistic_dictionary.update({f"urms {run_name}": urmse})
  statistic_dictionary.update({f"nrsd {run_name}": nrsd})
  statistic_dictionary.update({f"fac2 {run_name}": fac2})
  statistic_dictionary.update({f"mfbs {run_name}": mfbs})
  statistic_dictionary.update({f"aqpi {run_name}": aqpi})

  return statistics_string


class HistPlot:

    def __init__(self, ax_size, figure_size=(9,5)):
        """
        Initializes the histogram plot with a specified figure size.

        Args:
            ax_size: Size of the axes.
            figure_size (tuple, optional): Size of the figure (width, height) in inches.
                Defaults to (9, 5).
        """

        self.fig  = figure(figsize = (9,5))
        self.ax   = self.fig.add_axes(ax_size)
        self.observed_data = []                     # Stores observed values (added only once)
        self.model_data    = []                     # Stores model data

    def add_data(self, observed_values, model_data):
        """
        Adds observation and model values to the plot data.

        Args:
            observed_values (list): List of observed values.
            model_data (list): List of model values.
        """

        if len(self.observed_data) == 0:            # Add observed values only if empty
            self.observed_data = [np.array(observed_values)]
        self.model_data += [model_data]

    def finalize(self, site_name, run_names, colors):
        """
        Finalizes the histogram plot with legend, binning, density estimates,
        and visual elements.

        Args:
            site_name (str): Name of the site.
            run_names (list): List of names for each data run (observation and day).
            colors (list): List of colors for each data run.
        """

        # Combine observed and model data for plotting
        all_data = self.observed_data + self.model_data

        # Find minimum and maximum values for combined data
        hf = np.array(all_data).flatten()
        data_min, data_max = hf.min(), hf.max()                #havg = hf.mean()

        # Create bins
        num_bins = hist     #20  # Adjust this parameter as needed
        bins = np.linspace(data_min, data_max, num_bins)

        # Combine observation and day labels with site name
        labels = ([site_name] + run_names)

        obscolors = (['Red'] + colors)

        # Plot histograms
        hist_counts, bin_edges, patches = self.ax.hist(
            all_data,
            bins=bins,
            histtype='bar',
            label=labels,
            color=obscolors,
            linewidth=2
        )

        # Adjust legend placement
        self.ax.legend(
            bbox_to_anchor=(0, 1.02, 1, 0.1),
            loc=3,
            ncol=7,
            mode="expand",
            borderaxespad=0
        )

        # Generate x-axis values for density estimates
        dense_x = np.linspace(data_min, data_max, 1000)

        # Plot density estimates for each data run
        for i, data_chunk in enumerate(all_data):
            mu, std = norm.fit(data_chunk)  # Fit normal distribution
            density = norm.pdf(dense_x, mu, std)
            # Normalize and scale density for accurate representation
            density = density / density.sum() * data_chunk.sum() * (len(dense_x) / len(bins))
            self.ax.plot(dense_x, density, '--', color=obscolors[i], linewidth=1, label=labels[i])

        # Add labels and title
        #self.ax.set_xlabel("[Hg$^0$]")
        #self.ax.set_ylabel("Frequency")
        #self.fig.subtitle(f"Histogram of {site_name} Data")

        return self.fig


def PlotSite(site_data):
    """
    Plots time series data for a single site and various model runs.

    Args:
        site_data (tuple): A tuple containing data for the site:
            * site_name (str): Name of the site
            * obs_dates (np.ndarray): Dates of observations (datetime objects)
            * obs_values (np.ndarray): Observed values
            * obs_dates_hr (np.ndarray, optional): Dates of hourly observations (datetime objects)
            * obs_values_hr (np.ndarray, optional): Hourly observed values
            * smtw_obs_values (np.ndarray, optional): Weekly smoothed observed values
            * smtm_obs_values (np.ndarray, optional): Monthly smoothed observed values
            * obs_diurnal (np.ndarray, optional): Diurnal observed values
            * model_data (list): List containing model data for each run (list format)
                * time (np.ndarray): Time series for the model run
                * days (np.ndarray): Modelled daily values
                * week (np.ndarray): Modelled weekly values
                * mnth (np.ndarray): Modelled monthly values
                * diur (np.ndarray): Modelled diurnal values
                * timh (np.ndarray): Modelled hourly time series
                * hour (np.ndarray): Modelled hourly values
            * pdf_name (str): Name of the output PDF file

    Returns:
        dict: A dictionary containing various statistics for each model run.
    """

    site, obs_dat, obs_val, obs_dat_hr, obs_val_hr, smtw_obs_val, smtm_obs_val, obs_diu, mods, pdf_name = site_data

    stat = {}

    pp = PdfPages(pdf_name, metadata=metadata)

    obs_val_msk = np.logical_and(~np.isnan(obs_val), datemsk)   # mask to remove NaNs and excluded months
    obs_val_num = obs_val[obs_val_msk]                          # remove NaNs
    obs_avg = obs_val_num.mean()

    mpl.pyplot.style.use('ggplot')                              # Plotting style for matplotlib

    fig = figure(figsize = (10,5))
    ax  = fig.add_axes(ax_size)

    if args.statplot:
        statplot = ScatterPlotWithRegression()

    if diurplot:
        diurnal_plot = DiurnalPlot()
        diurnal_plot.plot_data(obs_diu, dih, dim, 'Obs', color='Red', alpha=0.85 if args.trans else 1, markersize=7, linewidth=2.5)

    if args.taylor:
        taylor_fig = figure(figsize = (9,6))
        mtd = ModTaylorDiagram(taylor_fig)

    if args.solar:
        solar_fig = figure(figsize = (9,6))

    if args.biascor:
        biascor_fig = figure(figsize = (9,6))

    if args.solar or args.biascor:
        obsvals = []
        modvals = []

    if hist>0:
        histplot = HistPlot(ax_size=ax_size)

    if args.daily:
        obsl = ax.plot_date(obs_dat[obs_val_msk] , obs_val[obs_val_msk], '.-', color='r', antialiased=True, alpha=0.99 if args.trans else 1, lw=2, markersize=2.0)
    if args.weekly:
        obsm = ax.plot_date(obs_dat              , smtw_obs_val        , '-' , color='r', antialiased=True, alpha=0.85 if args.trans else 1, lw=3)
    if args.monthly:
        obsm = ax.plot_date(obs_dat              , smtm_obs_val        , '-' , color='r', antialiased=True, alpha=0.85 if args.trans else 1, lw=3)

    modl = []
    modm = []
    modo = []
    modmin, modmax = 10, 0

    if args.stats:
        dtp = 0.025
        tp1 = 0.030 + len(runs)*dtp

    # Loop over model runs
    for (time, days, week, mnth, diur, timh, hour), run, color, runsname in zip(mods, runs, colors, runsnames):

        if len(days)>0 and np.abs(days[~np.isnan(days)].mean()-1.0)>0.0001:                               # do not plot model data which is outside domain

            if args.daily:
                yy = np.where(days>0, days, NaN)
                modmin = min(modmin, min(yy))
                modmax = max(modmax, max(yy))
                modl += ax.plot_date(time, yy, '.-', color=color, antialiased=True, alpha=0.75 if args.trans else 1, lw=2, markersize=2)

            if args.weekly:
                yz = np.where(week>0, week, NaN)
                modmin = min(modmin, min(yz))
                modmax = max(modmax, max(yz))
                modm += ax.plot_date(time, yz, '-', color=color, antialiased=True, alpha=0.85 if args.trans else 1, lw=2)

            if args.monthly:
                yx = np.where(mnth>0, mnth, NaN)
                modmin = min(modmin, min(yx))
                modmax = max(modmax, max(yx))
                modm += ax.plot_date(time, yx, '-', color=color, antialiased=True, alpha=0.85 if args.trans else 1, lw=2)

            valid = np.logical_and(np.logical_and(~np.isnan(obs_val), ~np.isnan(days)), datemsk)            # remove NaNs from both arrays, obs and mod as well as excluded months
            mod_avg = days[valid].mean()                                                                    # calculate model mean after remove NaNs and excluded months

            statplot.add_data(obs_val[valid], days[valid], alpha=0.85 if args.trans else 1, color=color) if args.statplot else None
            diurnal_plot.plot_data(diur, dih, dim, runsname, color, alpha=0.85 if args.trans else 1) if diurplot else None
            if args.taylor:
                nrsd, corrcoef, nesd, nbias = mtd.add_prediction(obs_val[valid], days[valid], runsname)
            if args.solar or args.biascor:
                obsvals += [obs_val[valid]]
                modvals += [days   [valid]]
            histplot.add_data(obs_val[valid], days[valid]) if hist>0 else None

            if args.stats:                                                                                  # calc and save stats
                st = calculate_statistics(obs_val_hr, hour, datemskhr, stat, run)
                fig.text(0.05, tp1, st, fontdict={'family':'monospace', 'color':color, 'weight':'normal', 'size':8 })
                tp1 -= dtp

    # format the ticks
    ax.xaxis.set_major_locator  (MonthsLoc)
    ax.xaxis.set_major_formatter(MonthsFmt)
    #ax.xaxis.set_minor_locator  (DaysLoc  )
    #ax.autoscale_view()

    ax.set_xlim(datemin, datemax)

    offset = mpl.transforms.ScaledTranslation(30/72, 0, fig.dpi_scale_trans)              # move x-axis labels to the center
    for label in ax.xaxis.get_majorticklabels():
        label.set_transform(label.get_transform() + offset)

    # Calculate and set the y-axis limits
    datamin = max( 0., float(int(min(min(obs_val_num), modmin) * 10) - 1) / 10) if args.ymin == None else args.ymin
    datamax = min(12., float(int(max(max(obs_val_num), modmax) * 10) + 1) / 10) if args.ymax == None else args.ymax
    ax.set_ylim(datamin, datamax)

    ax.grid(True)

    if args.daily:
        legend = (obsl + modl)
    if args.weekly or args.monthly:
        legend = (obsm + modm)
    if args.daily or args.weekly or args.monthly:
        sitename = f'{site} ({Lox[site][0]:.1f} {Lox[site][1]:.1f})' if args.lox and Lox[site] else site                                # Add site coordinatesif necessary
        ax.legend(legend, ([sitename] + runsnames), bbox_to_anchor=(0, 1.02, 1, 0.1), loc=3, ncol=7, mode="expand", borderaxespad=0.)   # The legend has site name for observations and run names
        pp.savefig(fig)

    pp.savefig(statplot.finalize(runsnames)) if args.statplot else None
    pp.savefig(diurnal_plot.finalize_plot(title=f"{sitename} Diurnal Plot", legend_args={'bbox_to_anchor':(0, 1.02, 1, 0.1), 'loc':3, 'ncol':7, 'mode':'expand', 'borderaxespad':0})) if diurplot else None
    if args.taylor:
        mtd.plot()
        pp.savefig(taylor_fig)
    pp.savefig(SolarDiagram(modvals, obsvals[0], modnames=runsnames, fig=solar_fig)) if args.solar else None
    if args.biascor:
        bias_correlation_diagram(modvals, obsvals[0], model_names=runsnames, fig=biascor_fig, color=colors)
        pp.savefig(biascor_fig)
    pp.savefig(histplot.finalize(site, runsnames, colors)) if hist>0 else None

    clf()
    fig.clear()
    plt.close()

    pp.close()

    return stat


#======================================================================= Program start =======================================================================

if args.parallel and len(sites)>1:
    # Create a Ray pool of worker processes
    print ("Read and process data for all sites in parallel")
    ray.init(num_cpus=40)       # Set the number of worker processes
    pool = Pool()
    data = pool.map(read_site, sites)
else:
    # No parallelization or small dataset, use sequential execution
    print ("Read and process data in serial mode")
    data = []
    for site in sites:
        print ("Process: "+site)
        data += [read_site(site)]

print ("Data was read and processed, start plotting.")

# Axes for plots
ax_size   = [0.04, 0.05, 0.95, 0.85]
MonthsLoc = MonthLocator()                                        # ticks for every month
MonthsFmt = DateFormatter('%b')
DaysLoc   = DayLocator()

# Create PDF file
pdf_name = "%4u" % args.year
for run in runs:
    pdf_name = pdf_name + '.' + run
pdf_name = pdf_name + "." + { True:'d', False:'' }[args.daily  ] + { True:'w', False:'' }[args.weekly  ] + \
                            { True:'m', False:'' }[args.monthly] + { True:'s', False:'' }[args.statplot] + \
                            { True:f'u{dim}{dih}', False:'' }[diurplot] + \
                            { True:'t', False:'' }[args.taylor ] + \
                            { True:'r', False:'' }[args.solar  ] + \
                            { True:'b', False:'' }[args.biascor] + \
                            { True:'h', False:'' }[hist!=0     ] + ".pdf"

inp = []
pdf_name_sites = []
for site, obs_dat, obs_val, obs_dat_hr, obs_val_hr, smtw_obs_val, smtm_obs_val, obs_diu, mods in data:
    pdf_name_site = pdf_name+'.'+site+'.pdf'
    inp += [(site, obs_dat, obs_val, obs_dat_hr, obs_val_hr, smtw_obs_val, smtm_obs_val, obs_diu, mods, pdf_name_site)]
    pdf_name_sites += [pdf_name_site]

if args.parallel and len(inp)>2:
    pool    = Pool()
    statall = pool.map(PlotSite, inp)
    pool.close()
    pool.join()
else:
    statall = []
    for i in inp:
        statall += [PlotSite(i)]


if args.stats:
    tts = ['Bias','Correlation','Slope','Intersect','StdErr','RMSE','URMSE','NormStDev','FAC2','MF Bias','AQPI']
    sts = ['bias','corr'       ,'slop' ,'intr'     ,'stdr'  ,'rmse','urms' ,'nrsd'     ,'fac2','mfbs'   ,'aqpi']

    for t, s in zip(tts, sts):
        print ("\n{:19s}".format(t), end='')
        for run in runs:
            print ("{:10s}".format(run[:10]), end='')
        print ('')
        for i, statone in enumerate(statall):                                                                # loop over sites
            print ("{:14s}: ".format(sites[i][:14]), end='')
            mx = -100; mn = 100
            for i,run in enumerate(runs):
                sa = statone[s+' '+run]
                if sa>mx:
                    mx = sa
                    ix = i
                if sa<mn:
                    mn = sa
                    im = i
            for i,run in enumerate(runs):
                if i==ix:
                    c='x'
                elif i==im:
                    c='*'
                else:
                    c=' '
                print("{: 9.4f}{:1s}".format(statone[s+' '+run], c),end='')
            print('')

    stat = {}
    stan = {}
    for run in runs:
        for st in sts:
            stat.update({st+' '+run:0.0})
            stan.update({st+' '+run:0.0})

    for statone in statall:                                                                                # loop over sites
        for s in statone:
            stat[s] += statone[s]
        for st in sts:
            statmin, statmax = 1000.0, -1000.0                                                                # find min&max for each statistics for given site
            for run in runs:
                statmin = min(statmin, statone[st+' '+run])
                statmax = max(statmax, statone[st+' '+run])
            for run in runs:
                statdif = statmax - statmin
                if statdif!=0:
                    stan[st+' '+run] += (statone[st+' '+run]-statmin)/statdif                        # normalize to 0...1

    for stai,styp in ((stat,'Aggr'),(stan,'Norm')):
        print ('\nModel:'+styp+'   ',end='')
        for st in sts:
            print ("{:7s}".format(st),end='')
        print ('')
        for run in runs:
            print ("{:12s}".format(run[:12]),end='')
            for st in sts:
                print ("{: 6.4f}".format(stai[st+' '+run]/len(sites)),end='')
            print ('')


    # Print combined statistics for all sites
    dat = dict()
    for i,site in enumerate(sites):
        obsd = data[i][2][datemsk]
        obsh = data[i][4][datemskhr]
        try:
            dat['ObsD'] = np.append( dat['ObsD'], obsd )
            dat['ObsH'] = np.append( dat['ObsH'], obsh )
        except:
            dat['ObsD'] = obsd
            dat['ObsH'] = obsh
        for r,m in enumerate(data[i][8]):
            modd = m[1][datemsk]
            modh = m[6][datemskhr]
            try:
                dat['ModD'+runs[r]] = np.append( dat['ModD'+runs[r]], modd )
                dat['ModH'+runs[r]] = np.append( dat['ModH'+runs[r]], modh )
            except:
                dat['ModD'+runs[r]] = modd
                dat['ModH'+runs[r]] = modh

    #np.savez_compressed(pdf_name+'.npz', dat=dat)

    for tim in ['D','H']:

      print ("\nAll sites "+{'D':'daily','H':'hourly'}[tim]+"")

      if args.taylor:
          pdf_name_taylor = pdf_name+'taylor'+tim+'.pdf'
          pdf_name_sites += [pdf_name_taylor]
          pp = PdfPages(pdf_name_taylor)
          taylor_fig = figure(figsize = (9,6))
          mtd = ModTaylorDiagram(taylor_fig)

      for i,r in enumerate(runs):
          obs = dat['Obs'+tim]
          mod = dat['Mod'+tim+r]
          valid = np.logical_and( ~np.isnan(obs), ~np.isnan(mod) )
          obs = obs[valid]
          mod = mod[valid]
          slope, intercept, corr, p_value, std_err = linregress(obs, mod)
          bias    = mod.mean() - obs.mean()
          rmse    = RMSE ( mod, obs )
          urmse   = URMSE( mod, obs )
          nrsd    = NSTD ( mod, obs )
          fac2    = FAC2 ( mod, obs )
          mfbs    = MFB  ( mod, obs )
          aqpi    = (fac2 + corr + (1.0-np.abs(0.5*mfbs)))/3.0
          print ("{:12s}: Bias={:+6.3f}  R={:6.4f}  mod={:5.2f}*obs{:+5.2f}  n={:6g}  RMSE={:4.3f}  URMSE={:4.3f}  NRSD={:4.3f}  FAC2={:4.3f}  MFB={:6.3f}  AQPI={:4.3f}".format(
                  r[:12], bias,          corr,      slope,  intercept,  obs.shape[0], rmse,        urmse,         nrsd,         fac2,         mfbs,        aqpi ) )

          if args.taylor:
              nrsd, corrcoef, nesd, nbias = mtd.add_prediction(obs, mod, runs[i], '%u'%(i+1))

      if args.taylor:
          mtd.plot()
          pp.savefig(taylor_fig)
          pp.close()


# Collect all PDFs into one
writer = PdfWriter(version="1.7", compress=pdfrw_compress)
info = None
size0 = 0
for inpfn in pdf_name_sites:
    reader = PdfReader(inpfn)
    if info == None:
        info = reader.Info
    writer.addpages(reader.pages)
    size0 += os.stat(inpfn).st_size
    os.remove(inpfn)
writer.trailer.Info = info
writer.write(pdf_name)

# Compress PDF file if cpdf is available
size1 = os.stat(pdf_name).st_size
if cpdf:
    pdf_name_comp = pdf_name[:-4]+'.c.pdf'
    
    #status = subprocess.call('cpdf -squeeze '+pdf_name+' -o '+pdf_name_comp+' &>/dev/null', shell=True)
    process = subprocess.Popen(['cpdf', '-squeeze', pdf_name, '-o', pdf_name_comp], stdout=subprocess.PIPE, stderr=subprocess.PIPE)

    # Wait for the process to finish
    output, error = process.communicate()

    if process.returncode == 0:
        print("cpdf completed successfully!")
    else:
        print(f"cpdf failed with error: {error.decode()}")

    size2 = os.stat(pdf_name_comp).st_size
    if os.path.isfile(pdf_name_comp):
        os.remove(pdf_name)
        os.rename(pdf_name_comp, pdf_name)
    print(f"\n{pdf_name} created. Size: {size0} => {size1} => {size2} :: {size2*100//size0:3d}%")
    #pdf_name = pdf_name_comp
else:
    print(f"\n{pdf_name} created. Size: {size0} => {size1} :: {size1*100//size0:3d}%")


# Create png and pptx files if necessary
if dpi>0:

    filenames = []

    seq = convert_from_path(pdf_name, dpi, fmt='png')
    for Page, img in enumerate(seq):
        filename = f'{pdf_name[:-3]:s}{Page+1:03d}.png'
        print(f"Convert PDF to PNG [{filename}] :: ", end='')
        img = img.convert('RGB')
        if pngclr<=256:
            img = img.convert('P', colors=pngclr, palette=Image.ADAPTIVE, dither=0)                       # Convert to a palette          Image.FLOYDSTEINBERG
            #img = img.quantize(pngclr, dither=0)
        img.save(filename)

        if args.pngopt:
            size0 = os.stat(filename).st_size
            time0 = timeit.default_timer()
            if shutil.which('optipng'):
                process = subprocess.Popen([shutil.which('optipng'), '-o1', '-strip', 'all', '-quiet', filename], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                output, error = process.communicate()
            time1 = timeit.default_timer()
            size1 = os.stat(filename).st_size
            #if shutil.which('pngout'):
            #    #status = subprocess.call('pngout -f5 -s1 -b128 -c3 -d8 '+filename+' &>/dev/null', shell=True)
            #    process = subprocess.Popen([shutil.which('pngout'), '-f5', '-s1', '-b128', '-c3', '-d8', filename], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            #    output, error = process.communicate()
            time2 = timeit.default_timer()
            size2 = os.stat(filename).st_size
            if shutil.which('advpng'):
                process = subprocess.Popen([shutil.which('advpng'), '-z', '-3', '-q', filename], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                output, error = process.communicate()
            time3 = timeit.default_timer()
            size3 = os.stat(filename).st_size
            print(f'{size0:6d} => {size1:6d} => {size2:6d} => {size3:6d} :: {size3*100//size0:3d}% ::: {time1-time0:4.1f} : {time2-time1:4.1f} : {time3-time2:4.1f}')
        else:
            print(f'{os.stat(filename).st_size:6d}')

        filenames += [filename]

    # Create pptx file
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for filename in filenames:
        slide = prs.slides.add_slide(blank_slide_layout)
        left, top, width = Cm(1), Cm(1.5), Cm(23.4)
        pic = slide.shapes.add_picture(filename, left, top, width=width)
        os.remove(filename)

    pptx_file = pdf_name[:-3]+'pptx'
    prs.save(pptx_file)
    print (pptx_file+' was created.')
